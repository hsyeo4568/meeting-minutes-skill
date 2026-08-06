#!/usr/bin/env python3
"""Phase 1.5 materials-digest floor: meeting materials -> one digest .md each.

This is the *floor* the pipeline falls back to when no configured handler skill
is installed (see `references/engine/pipeline.md` §1.5). It exists because a
live run on 2026-07-29 re-derived this logic ad hoc and reproduced two defects
worth locking down:

  - a deck digest built from extracted evidence alone is number soup: the
    deterministic evidence contract carries numbers-with-ids, while the slide's
    prose comes from a separate per-slide split;
  - that split is keyed 1-based, so a 0-based lookup hands every slide its
    neighbour's text and nothing complains.

Deterministic only. Rendered-page vision reading (`materials.deep_read`) is a
separate, costed decision — this module flags what it could not read instead of
implying the slide was understood.

Usage:
  python scripts/materials_digest.py <file|folder> [--out DIR]
"""
from __future__ import annotations

import argparse
import pathlib
import sys

if sys.stdout is not None:                    # scheduled runs have no console
    sys.stdout.reconfigure(encoding="utf-8")

# ext -> handler id. Unknown extensions return None so the caller can surface
# them as skipped_unsupported rather than digesting them into silence.
_ROUTES = {
    ".pptx": "pptx",
    ".pdf": "pdf",
    ".docx": "docx",
    ".html": "docx",
    ".xlsx": "table",
    ".csv": "table",
    ".txt": "text",
    ".md": "text",
}

_MAX_EVIDENCE_PER_GROUP = 40


def route(ext: str) -> str | None:
    """Handler id for a file extension, or None when unsupported."""
    return _ROUTES.get(ext.lower())


# --- pptx ------------------------------------------------------------------

def _ev_text(e) -> str:
    return f"{e.value}{getattr(e, 'unit', '') or ''}"


def render_pptx_digest(name: str, units, prose: dict) -> str:
    """Compose one deck digest from per-slide units + the 1-based prose split.

    `units` are slide extracts (index, hidden, evidence, images); `prose` maps
    **1-based** slide number -> slide text. Both halves are required: prose
    alone loses the figures, evidence alone loses the meaning.
    """
    lines = [f"# {name} — {len(units)} slides (deterministic extract)"]
    for u in units:
        body = (prose.get(u.index + 1) or "").strip()
        evidence = list(u.evidence)
        images = list(u.images)
        lines.append(f"\n## slide {u.index} (hidden={u.hidden}) "
                     f"ev={len(evidence)} img={len(images)}")
        lines.append(body if body else "_본문 텍스트 없음 — 도식/이미지 슬라이드 가능성_")
        by_loc = {"shape": [], "chart": [], "table": [], "note": []}
        for e in evidence:
            for key in by_loc:
                if str(e.loc).startswith(key):
                    by_loc[key].append(e)
                    break
        for key, items in by_loc.items():
            if not items:
                continue
            shown = items[:_MAX_EVIDENCE_PER_GROUP]
            rendered = ", ".join(f"{_ev_text(e)} [{e.eid}]" for e in shown)
            more = "" if len(items) == len(shown) else f" (+{len(items) - len(shown)} more)"
            lines.append(f"- {key}: {rendered}{more}")
        if images:
            ids = ", ".join(i.image_id for i in images)
            lines.append(f"- images: {ids} — 배치·도식 의미 미판독 (deep_read 별도 결정)")
    return "\n".join(lines)


def _exec_module(name: str, path: pathlib.Path):
    """Load a module from a file path.

    The module is registered in ``sys.modules`` *before* execution: a module
    defining ``@dataclass`` reads ``sys.modules[cls.__module__]`` while the
    decorator runs and dies with a bare AttributeError otherwise. That failure
    is invisible here — the caller just degrades to the floor extractor.
    """
    import importlib.util

    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    try:
        spec.loader.exec_module(module)
    except Exception:
        sys.modules.pop(name, None)
        raise
    return module


def _load_deck_engine(hub: pathlib.Path | None = None):
    """Locate an installed sibling deck-distiller engine, or None.

    Resolved by capability, not by name: a sibling skill exposing an
    ``extract_deck`` + ``split_markitdown`` pair. Without this lookup the
    'prefer the distiller' branch would never fire outside that skill's own
    test run — decorative preference, floor quality forever.
    """
    hub = hub or pathlib.Path(__file__).resolve().parent.parent.parent
    for sibling in sorted(p for p in hub.iterdir() if p.is_dir()):
        engine = sibling / "references" / "engine" / "extract.py"
        if not (sibling / "SKILL.md").exists() or not engine.exists():
            continue
        if str(sibling) not in sys.path:
            sys.path.insert(0, str(sibling))          # `references` is a namespace pkg under ROOT
        try:
            module = _exec_module(f"_deck_engine_{sibling.name}".replace("-", "_"), engine)
        except Exception:
            continue
        if hasattr(module, "extract_deck") and hasattr(module, "split_markitdown"):
            return module
    return None


def _digest_pptx(path: pathlib.Path) -> str:
    """Prefer an installed deck-distiller engine; fall back to python-pptx."""
    engine = _load_deck_engine()
    if engine is not None:
        units = engine.extract_deck(str(path))
        return render_pptx_digest(path.name, units, engine.split_markitdown(str(path)))
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = [f"# {path.name} — {len(prs.slides)} slides (python-pptx floor)"]
    for i, slide in enumerate(prs.slides):
        texts = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        pictures = [sh.shape_type for sh in slide.shapes if sh.shape_type == 13]
        lines.append(f"\n## slide {i}")
        lines.append("\n".join(texts) if texts else "_본문 텍스트 없음_")
        if pictures:
            lines.append(f"- images: {len(pictures)} — 도식 의미 미판독 (deep_read 별도 결정)")
    return "\n".join(lines)


# --- other formats ---------------------------------------------------------

def _digest_docx(path: pathlib.Path) -> str:
    from markitdown import MarkItDown
    return f"# {path.name} (markitdown)\n\n{MarkItDown().convert(str(path)).text_content}"


def _digest_pdf(path: pathlib.Path) -> str:
    import fitz
    doc = fitz.open(str(path))
    out = [f"# {path.name} — {doc.page_count} pages (text layer)"]
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        out.append(f"\n## p{i + 1}\n" + (text or "_텍스트 레이어 없음 — 스캔본 가능성_"))
    return "\n".join(out)


def _digest_table(path: pathlib.Path) -> str:
    """Schema first: sheet/column names before any values (whitespace + case bite)."""
    import pandas as pd
    if path.suffix.lower() == ".csv":
        frames = {"csv": pd.read_csv(path)}
    else:
        frames = pd.read_excel(path, sheet_name=None)
    out = [f"# {path.name} — {len(frames)} sheet(s)"]
    for sheet, df in frames.items():
        out.append(f"\n## {sheet} — {len(df)} rows")
        out.append("- columns: " + ", ".join(map(str, df.columns)))
        out.append(df.head(5).to_string(index=False))
    return "\n".join(out)


def _digest_text(path: pathlib.Path) -> str:
    return f"# {path.name}\n\n{path.read_text(encoding='utf-8', errors='replace')}"


_HANDLERS = {
    "pptx": _digest_pptx,
    "pdf": _digest_pdf,
    "docx": _digest_docx,
    "table": _digest_table,
    "text": _digest_text,
}


def digest_file(path: pathlib.Path, outdir: pathlib.Path) -> tuple[str, str]:
    """Digest one material. Returns (status, detail) — never raises for an
    unsupported extension; the caller reports it."""
    handler = route(path.suffix)
    if handler is None:
        return "skipped_unsupported", path.suffix
    text = _HANDLERS[handler](path)
    dest = outdir / f"{path.stem}.digest.md"
    dest.write_text(text, encoding="utf-8")
    return "ok", f"{handler} -> {dest.name} ({len(text)} chars)"


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("target", help="material file or folder of materials")
    ap.add_argument("--out", default=None, help="digest output dir (default: <target>/_digest)")
    args = ap.parse_args()

    target = pathlib.Path(args.target)
    files = sorted(p for p in target.iterdir() if p.is_file()) if target.is_dir() else [target]
    outdir = pathlib.Path(args.out) if args.out else (
        target if target.is_dir() else target.parent) / "_digest"
    outdir.mkdir(parents=True, exist_ok=True)

    failures = 0
    for path in files:
        try:
            status, detail = digest_file(path, outdir)
        except Exception as exc:                       # a broken material must not kill the run
            status, detail = "failed", f"{type(exc).__name__}: {exc}"
            failures += 1
        print(f"{status:<20} {path.name:<50} {detail}")
    print(f"\ndigests: {outdir}")
    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(main())
